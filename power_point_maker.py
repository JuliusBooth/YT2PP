from pptx import Presentation
from pptx.dml.color import RGBColor
from image_search import get_image_from_keyword
from pptx.util import Inches
from pptx.oxml import parse_xml

def add_image_to_slide(prs, slide, keyword, slide_color):
    image_path = get_image_from_keyword(keyword)
    if image_path:
        height = prs.slide_height
        width = prs.slide_width
        left = top = Inches(0)
        picture = slide.shapes.add_picture(image_path, left, top, width, height)
        slide.shapes._spTree.insert(2, picture._element)
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*slide_color)

def add_outline(p):
    p_xml_obj = p._p
    new_xml = """
        <a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:defRPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" b="1">
                <a:ln w="19000" cap="flat">
                    <a:solidFill>
                        <a:srgbClr val="000000"/>
                    </a:solidFill>
                    <a:prstDash val="solid"/>
                    <a:miter lim="400000"/>
                </a:ln>
                <a:solidFill>
                    <a:srgbClr val="FFFFFF"/>
                </a:solidFill>
            </a:defRPr>
        </a:pPr>
    """
    new_element = parse_xml(new_xml)
    p_xml_obj.append(new_element)

def add_bullets_to_slide(bullet_points, slide):
    content = slide.placeholders[1]
    for point in bullet_points:
        p = content.text_frame.add_paragraph()
        add_outline(p)
        run = p.add_run()
        run.text = point

def add_title_to_slide(slide_title, slide):
    title = slide.shapes.title
    p = title.text_frame.paragraphs[0]
    add_outline(p)
    run = p.add_run()
    run.text = slide_title

def create_presentation(presentation_title, slides_data):
    prs = Presentation()
    for slide_obj in slides_data:
        slide_title = slide_obj["Title"]
        bullet_points = slide_obj["Bullets"]
        slide_color = slide_obj["Slide Color"]
        image_keyword = slide_obj["Image"]
        slide_layout = prs.slide_layouts[1]  # 1 is the index of a slide layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        add_image_to_slide(prs, slide, image_keyword, slide_color)
        add_title_to_slide(slide_title, slide)
        add_bullets_to_slide(bullet_points, slide)
    prs.save(f'power_points/{presentation_title}.pptx')


if __name__ == "__main__":
    data = [{'Title': 'Starting Your Presentation', 'Bullets': ['Create a specific working title and main takeaways', 'Consider audience knowledge and expectations', 'Formulate content as a narrative and use compelling stories', 'Collect credible data and examples to support your argument'], 'Image': 'Presentation', 'Slide Color': (106, 159, 219)}, {'Title': 'Designing Your PowerPoint', 'Bullets': ['Choose a color scheme that supports your content and is easily legible', 'Use pre-designed templates as a starting point for background design', 'Format content with concise language and main phrases', 'Begin and end with engaging points of interest'], 'Image': 'PowerPoint Design', 'Slide Color': (219, 106, 159)}]
    create_presentation("how to make a good powerpoint",data)